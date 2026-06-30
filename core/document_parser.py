"""
统一文档解析器
支持格式：pptx / pdf / docx / doc / md / txt → 统一输出纯文本
"""
import os
from pathlib import Path
from typing import Optional


class DocumentParser:
    """将各种格式的文档解析为纯文本"""

    # ----------------------------------------------------------------
    # 公有入口
    # ----------------------------------------------------------------

    @staticmethod
    def parse(file_path: str) -> str:
        """
        自动识别文件类型并解析为纯文本
        Args:
            file_path: 文件路径
        Returns:
            解析后的纯文本字符串
        """
        ext = Path(file_path).suffix.lower()
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        parsers = {
            ".pptx": DocumentParser._parse_pptx,
            ".ppt": DocumentParser._parse_ppt_old,
            ".pdf": DocumentParser._parse_pdf,
            ".docx": DocumentParser._parse_docx,
            ".doc": DocumentParser._parse_doc_old,
            ".md": DocumentParser._parse_markdown,
            ".txt": DocumentParser._parse_txt,
        }

        parser = parsers.get(ext)
        if parser is None:
            raise ValueError(f"不支持的文件格式: {ext}。支持: {list(parsers.keys())}")

        return parser(file_path)

    @staticmethod
    def parse_with_metadata(file_path: str) -> dict:
        """
        解析文件，同时返回元信息
        Returns:
            {"text": str, "file_name": str, "pages": int, "format": str}
        """
        text = DocumentParser.parse(file_path)
        return {
            "text": text,
            "file_name": Path(file_path).name,
            "format": Path(file_path).suffix.lower(),
            "char_count": len(text),
        }

    # ----------------------------------------------------------------
    # 各格式私有解析方法
    # ----------------------------------------------------------------

    @staticmethod
    def _parse_pptx(file_path: str) -> str:
        """解析 .pptx 文件，提取每页幻灯片的文本"""
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            lines.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        lines.append(" | ".join(cells))
            if lines:
                slides_text.append(f"【幻灯片 {i}】\n" + "\n".join(lines))

        return "\n\n".join(slides_text) if slides_text else "（PPT 无文本内容）"

    @staticmethod
    def _parse_ppt_old(file_path: str) -> str:
        """解析旧版 .ppt 文件。尝试转成 PDF 再读；失败则返回提示"""
        # 旧版 .ppt 是二进制格式，没有原生 Python 库能很好解析
        # 实际项目建议将 .ppt 另存为 .pptx 或 .pdf 后再导入
        raise ValueError(
            "不支持旧版 .ppt 格式。请用 PowerPoint 打开后另存为 .pptx 或 .pdf 格式再导入。"
        )

    @staticmethod
    def _parse_pdf(file_path: str) -> str:
        """解析 .pdf 文件，提取所有页面文本"""
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        pages_text = []
        for i, page in enumerate(doc, 1):
            text = page.get_text().strip()
            if text:
                pages_text.append(f"【第 {i} 页】\n{text}")

        doc.close()
        return "\n\n".join(pages_text) if pages_text else "（PDF 无可提取文本，可能是扫描件图片）"

    @staticmethod
    def _parse_docx(file_path: str) -> str:
        """解析 .docx 文件，提取段落和表格"""
        from docx import Document

        doc = Document(file_path)
        parts = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # 识别标题样式
                if para.style.name.startswith("Heading"):
                    level = para.style.name.split()[-1]
                    prefix = "#" * int(level) if level.isdigit() else "##"
                    parts.append(f"\n{prefix} {text}\n")
                else:
                    parts.append(text)

        for i, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            parts.append(f"\n【表格 {i+1}】\n" + "\n".join(rows))

        return "\n\n".join(parts)

    @staticmethod
    def _parse_doc_old(file_path: str) -> str:
        """
        解析旧版 .doc 文件。
        旧版 .doc 是二进制格式，Python 生态中没有能 100% 可靠解析的库。
        这里尽力而为，若失败则给出明确的转换指引。
        """
        # 策略优先级：
        # ① 尝试 antiword（如果系统装了）— 最可靠
        # ② 尝试 python-pptx 兼容读取 + olefile 提取
        # ③ 都失败则报错，引导用户转换

        # --- 策略 ①: antiword ---
        try:
            import subprocess
            result = subprocess.run(
                ["antiword", "-m", "UTF-8", file_path],
                capture_output=True, timeout=30,
            )
            text = result.stdout.decode("utf-8", errors="ignore").strip()
            text = result.stdout.strip()
            if text and DocumentParser._is_valid_text(text):
                return text
        except Exception:
            pass

        # --- 策略 ②: olefile 提取文本流 ---
        try:
            import olefile
            ole = olefile.OleFileIO(file_path)

            text_parts = []
            # WordDocument 主文本流
            for stream_name in ["WordDocument", "1Table", "0Table"]:
                if ole.exists(stream_name):
                    try:
                        data = ole.openstream(stream_name).read()
                        # 尝试提取其中的可读文本片段
                        text = data.decode("utf-8", errors="ignore")
                        text = "".join(c for c in text if c.isprintable() or c in "\n\r\t")
                        if DocumentParser._is_valid_text(text):
                            text_parts.append(text)
                    except Exception:
                        pass
            ole.close()

            combined = "\n".join(text_parts)
            if combined and DocumentParser._is_valid_text(combined):
                return combined
        except Exception:
            pass

        # --- 全部失败 ---
        raise ValueError(
            "无法解析旧版 .doc 文件。该格式是 Microsoft Word 97-2003 私有二进制格式，"
            "Python 生态无可靠解析库。\n"
            "请用 Word / WPS 打开后另存为 .docx 或 .pdf 格式再导入。"
        )

    @staticmethod
    def _is_valid_text(text: str, min_ratio: float = 0.3) -> bool:
        """
        检测文本是否为有效内容（非乱码）
        通过统计 CJK 字符 + 字母 + 数字的占比来判断
        """
        if not text or len(text) < 50:
            return False

        valid_count = 0
        for ch in text:
            cp = ord(ch)
            # CJK 统一汉字 + 扩展区
            if 0x4E00 <= cp <= 0x9FFF:      valid_count += 1
            elif 0x3400 <= cp <= 0x4DBF:     valid_count += 1
            elif 0x20000 <= cp <= 0x2A6DF:   valid_count += 1
            # 拉丁字母 + 数字
            elif ord('a') <= cp <= ord('z'): valid_count += 1
            elif ord('A') <= cp <= ord('Z'): valid_count += 1
            elif ord('0') <= cp <= ord('9'): valid_count += 1
            # 常见标点
            elif cp in (ord(c) for c in "，。；：！？、""''（）【】《》…—\n\r\t .:;!?()[]"):
                valid_count += 1

        ratio = valid_count / max(len(text), 1)
        return ratio >= min_ratio


    @staticmethod
    def _parse_markdown(file_path: str) -> str:
        """解析 .md 文件，直接读取（保留原始格式供后续处理）"""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    @staticmethod
    def _parse_txt(file_path: str) -> str:
        """解析 .txt 文件"""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()


# ----------------------------------------------------------------
# 便捷函数
# ----------------------------------------------------------------

def parse_document(file_path: str) -> str:
    """便捷函数：解析文档 → 纯文本"""
    return DocumentParser.parse(file_path)


def batch_parse(directory: str, extensions: Optional[list] = None) -> list[dict]:
    """
    批量解析目录下所有支持的文档
    Args:
        directory: 目录路径
        extensions: 限定后缀列表，如 ['.pptx', '.pdf']
    Returns:
        [{"text": ..., "file_name": ..., "format": ...}, ...]
    """
    if extensions is None:
        extensions = [".pptx", ".pdf", ".docx", ".doc", ".md", ".txt"]

    results = []
    for root, _, files in os.walk(directory):
        for fname in files:
            if Path(fname).suffix.lower() in extensions:
                file_path = os.path.join(root, fname)
                try:
                    result = DocumentParser.parse_with_metadata(file_path)
                    results.append(result)
                except Exception as e:
                    print(f"[跳过] {fname}: {e}")

    return results
